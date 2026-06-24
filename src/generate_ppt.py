# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 색상 정의 (Premium Dark Mode Palette)
BG_COLOR = RGBColor(15, 23, 42)        # Slate 900 (매우 어두운 네이비 그레이)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Slate 800 (카드 배경)
TEXT_WHITE = RGBColor(255, 255, 255)   # 타이틀 흰색
TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400 (설명/본문 회색)
ACCENT_TEAL = RGBColor(20, 184, 166)   # Teal 500 (핵심 하이라이트)
ACCENT_BLUE = RGBColor(59, 130, 246)   # Blue 500 (서브 하이라이트)
ACCENT_ORANGE = RGBColor(249, 115, 22) # Orange 500 (경고/통계적 오류 하이라이트)

FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Arial"

def create_slide_with_bg(prs):
    """어두운 배경이 적용된 슬라이드 생성"""
    blank_layout = prs.slide_layouts[6] # 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)
    
    # 배경색 지정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    return slide

def add_header(slide, title_text, category_text="PORTFOLIO PRESENTATION"):
    """상단 공통 헤더 추가"""
    # 카테고리 텍스트
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_BODY
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_TEAL
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def add_card(slide, left, top, width, height, title_text, text_lines):
    """SaaS 스타일 내용 카드 추가"""
    # 카드 모양 (둥근 직사각형)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG_COLOR
    shape.line.color.rgb = RGBColor(71, 85, 105) # Slate 600 경계선
    shape.line.width = Pt(1)
    
    # 카드 텍스트 프레임 생성
    tb = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + 0.25), 
        Inches(width - 0.5), Inches(height - 0.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # 카드 타이틀
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = ACCENT_TEAL
    p_title.space_after = Pt(10)
    
    # 카드 본문 라인 추가
    for line in text_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)
        
        # 특정 키워드가 포함될 경우 폰트 색상 및 굵기 하이라이트
        if "★" in line or "**" in line:
            clean_text = line.replace("★", "").replace("**", "")
            p.text = clean_text
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
        elif "핵심 역할" in line or "분석 기법" in line or "성과" in line or "의사결정" in line:
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE

def add_image_or_placeholder(slide, image_name, left, top, width, height):
    """지정된 이미지를 배치하고, 없을 경우 대안 플레이스홀더 카드 렌더링"""
    brain_dir = r"C:\Users\gnsl1\.gemini\antigravity-ide\brain\88edeb10-dad7-4721-93ed-5bdbacbab2f0"
    image_path = os.path.join(brain_dir, image_name)
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
            return
        except Exception as e:
            print(f"Failed to add image {image_name}: {e}")
            
    # Fallback placeholder card (이미지가 없을 경우)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG_COLOR
    shape.line.color.rgb = ACCENT_TEAL
    shape.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + height/2 - 0.5), Inches(width - 0.6), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[Streamlit Dashboard Image: {image_name}]"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL

def add_zoom_slide(prs, slide_num, title_text, image_name, notes_script):
    """실시간 시뮬레이션 전체 화면 줌인(Zoom-in) 슬라이드 생성"""
    slide = create_slide_with_bg(prs)
    add_header(slide, f"{slide_num}. [상세 화면] {title_text}", "DASHBOARD ZOOM-IN")
    
    # 꽉 찬 전체 화면 이미지 배치 (가로 11.3인치 x 세로 5.0인치)
    add_image_or_placeholder(slide, image_name, 1.0, 1.7, 11.33, 5.0)
    
    # 발표자 노트 설정
    set_presenter_notes(slide, notes_script)

def set_presenter_notes(slide, notes_text):
    """슬라이드 노트(발표자 폰 메모) 설정"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def main():
    prs = Presentation()
    # 16:9 와이드스크린 치수 설정 (13.33인치 x 7.5인치)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 공통 레이아웃 좌표 설정 (세로폭 확장: top=1.5, height=5.2)
    CARD_LEFT = 0.8
    CARD_TOP = 1.5
    CARD_WIDTH = 5.6
    CARD_HEIGHT = 5.2
    
    IMG_LEFT = 6.9
    IMG_TOP = 1.5
    IMG_WIDTH = 5.6
    IMG_HEIGHT = 4.2
    
    # =========================================================================
    # SLIDE 1: Title Slide (표지)
    # =========================================================================
    slide1 = create_slide_with_bg(prs)
    
    # 좌측 장식 바 (Teal)
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_TEAL
    accent_bar.line.fill.background()
    
    # 타이틀 텍스트 박스
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # 카테고리 레이블
    p0 = tf.paragraphs[0]
    p0.text = "DATA ANALYTICS & DECISION PORTFOLIO"
    p0.font.name = FONT_BODY
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_TEAL
    p0.space_after = Pt(20)
    
    # 메인 타이틀
    p1 = tf.add_paragraph()
    p1.text = "비정형 iGaming 로그 분석을 통한\n플레이어 성향 분류 및 의사결정 시뮬레이터 구축"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(14)
    
    # 서브타이틀
    p2 = tf.add_paragraph()
    p2.text = "비정형 데이터 파이프라인(ETL)부터 통계 분석, 머신러닝, 실시간 대시보드 웹 배포까지"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(40)
    
    # 메타정보
    p3 = tf.add_paragraph()
    p3.text = "SK Planet 데이터 직무 PT 면접 발표자료 (5~10분 분량)\n발표자: 지원자"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT_BLUE
    p3.font.bold = True
    
    # 슬라이드 1 노트 설정
    set_presenter_notes(slide1, """[슬라이드 1 발표 대본]
"안녕하십니까, SK Planet 데이터 직무 지원자 [지원자 이름]입니다.
저는 오늘 '비정형 iGaming 로그 분석을 통한 플레이어 성향 분류 및 의사결정 시뮬레이터 구축' 프로젝트를 주제로 발표를 진행하겠습니다.

본 프로젝트는 게임 시스템이 남긴 정제되지 않은 문자열 로그에서 출발하여, 핵심 지표를 계산하고, 통계 분석과 머신러닝을 적용해 유저의 의사결정을 돕는 웹 플랫폼까지 완성해 낸 End-to-End 데이터 사이언스 파이프라인 프로젝트입니다. 발표를 시작하겠습니다."

[용어 리마인더]
- 비정형 데이터: 컴퓨터가 바로 계산할 수 없는 자유 줄글 형식의 로그 텍스트 파일.""")
    
    # =========================================================================
    # SLIDE 2: 01. 프로젝트 개요 및 플랫폼 활성 지표 요약 (Tab 1: Overview)
    # =========================================================================
    slide2 = create_slide_with_bg(prs)
    add_header(slide2, "01. 프로젝트 개요 및 플랫폼 활성 지표 요약", "OVERVIEW & BASELINE")
    
    add_card(
        slide2, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT, 
        "비즈니스 현황 요약 및 기준선(Hero)",
        [
            "■ 분석 배경 및 문제 정의",
            "  - iGaming 플랫폼의 핵심 과제: 유저 조기 파산 및 이탈(Churn) 방어",
            "  - 올바른 베팅 가이드라인 부재로 인한 유저들의 급격한 칩 손실 대응 필요",
            "",
            "■ 분석 데이터셋 규모",
            "  - 140개 비정형 텍스트 로그 파일 파싱 및 적재",
            "  - 총 7,993개 핸드(게임 판수), 932명 플레이어, 10만 행 이상의 액션 로그",
            "",
            "■ ★대시보드 Overview 화면 연동",
            "  - **전체 활성 지표 요약**: 총 게임 수, 플레이어 수의 실시간 조회 메트릭",
            "  - **Hero 본인 지표 진단**: Hero의 VPIP(자발적 참여율), PFR(선제 공격율), AF(공격성 수치) 및 승률을 기준선으로 로드해 전체와 대조"
        ]
    )
    add_image_or_placeholder(slide2, "tab1_overview_1781002564252.png", IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    
    set_presenter_notes(slide2, """[슬라이드 2 발표 대본]
"첫 번째로 프로젝트의 개요와 대시보드 메인 화면인 '핵심 지표 요약(Overview)' 탭에 대해 말씀드리겠습니다.

포커와 같은 iGaming 플랫폼의 최대 과제는 유저들이 초반에 잘못된 플레이 습관으로 빠르게 자산을 잃어 파산하고 서비스를 이탈(Churn)하는 것을 막는 일입니다. 이를 위해 분석의 첫 단계로, 140개 비정형 텍스트 파일로부터 총 7,993판의 게임과 932명의 유저 데이터를 정제하여 적재했습니다.

우측에 보이시는 대시보드 화면이 바로 이를 요약한 화면입니다. 상단에는 전체 활성 지표 요약 카드(총 게임 수, 플레이어 수 등)가 표시되며, 중앙에는 분석 대상 본인인 'Hero'의 개인 베팅 스타일 지표인 VPIP, PFR, AF 및 승률을 데이터베이스에서 실시간으로 계산해 띄워 줍니다. 이 'Hero'의 지표를 전체 932명 플레이어들의 평균치와 비교 분석하는 것이 이번 의사결정 분석의 든든한 기준선이 되었습니다."

[용어 리마인더]
- Hero: 분석의 대상이 되는 본인 플레이어를 칭하는 포커 도메인 용어.
- VPIP(자발적 참여율): 기본 참가비 외에 자기 의지로 베팅에 참여한 판의 비율.
- PFR(선제 공격율): 판의 시작 단계에서 먼저 베팅액을 올려 기선을 제압한 비율.
- AF(공격성 수치): 판의 후반 단계에서 상대의 베팅에 단순 콜로 따라가기보다 레이즈나 베팅으로 적극적인 공격을 가한 횟수의 비율.""")

    # =========================================================================
    # SLIDE 3: ZOOM-IN Slide (01. 핵심 지표 요약 줌인)
    # =========================================================================
    add_zoom_slide(
        prs, "01-1", "핵심 지표 요약 대시보드 상세", "tab1_overview_1781002564252.png",
        """[슬라이드 상세 대본]
"우측의 대시보드 화면을 전체 화면으로 확대해 자세히 짚어드리겠습니다.
화면 최상단을 보시면 '플랫폼 활성 지표 요약' 카드가 존재합니다. 총 게임 수 7,993회, 플레이어 932명 등의 핵심 정보가 표시되어 있습니다. 그 아래 영역에서는 Hero 본인의 VPIP, PFR, AF 및 승률을 데이터베이스에서 즉각적으로 연산하여 시각적으로 보여줍니다. 

이처럼 전체 유저 수와 평균 지표를 실시간으로 스캔할 수 있는 인프라가 갖추어짐에 따라, 개별 플레이어들의 세부 성향을 비교 분석할 수 있는 기반이 마련되었습니다."
"""
    )

    # =========================================================================
    # SLIDE 4: 02. 플레이어 행동 프로파일링 데이터 마트 구축 (Tab 2: Profiling)
    # =========================================================================
    slide3 = create_slide_with_bg(prs)
    add_header(slide3, "02. 플레이어 행동 프로파일링 데이터 마트 구축", "PLAYER PROFILING & DATA MART")
    
    add_card(
        slide3, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT, 
        "932명 데이터 정형화 및 다차원 분류 마트",
        [
            "■ 플레이어 행동 데이터 마트 설계",
            "  - 932명 전체 플레이어의 VPIP, PFR, AF, 승률 등 개별 통계 집계",
            "  - 닉네임 검색 및 성향 필터를 통한 실시간 SQL 데이터 필터링 연동",
            "",
            "■ ★VPIP vs PFR 행동 분포 매트릭스 시각화",
            "  - X축(VPIP - 자발적 참여율), Y축(PFR - 선제 공격율) 기준의 2차원 산점도",
            "  - 플레이어들을 Tight-Aggressive(TAG), Loose-Passive(LP) 등 4대 성향 코호트 평면 상에 매핑하여 시각적 분포 규명",
            "",
            "■ 본인의 기술적 인프라 역할",
            "  - **SQLite 로컬 DB & PostgreSQL 엔터프라이즈 DB 이중화 설계**",
            "  - 문법 및 바인딩 기호 차이를 극복하는 추상화 커넥터 직접 구현"
        ]
    )
    add_image_or_placeholder(slide3, "tab2_profiling_1781002574421.png", IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    
    set_presenter_notes(slide3, """[슬라이드 3 발표 대본]
"두 번째로 '플레이어 프로파일링(Profiling)' 탭과 이를 뒷받침하는 데이터 마트 인프라에 대해 설명해 드리겠습니다.

파싱된 데이터를 기반으로, 저는 932명 플레이어 개개인의 VPIP, PFR, AF, 승률 수치를 전처리 및 집계하여 플레이어 데이터 마트 테이블을 구축했습니다. 그리고 사용자가 닉네임이나 플레이 스타일로 유저 목록을 실시간 필터링할 수 있는 SQL 조회 로직을 적용했습니다.

우측 화면의 하단을 보시면 'VPIP vs PFR 행동 분포 차트' 산점도가 렌더링되어 있습니다. X축은 게임 참견율인 VPIP이고, Y축은 선제공격율인 PFR입니다. 이 두 축을 기준으로 유저들의 행동 포인트가 평면 상에 흩어지며, 타이트-어그레시브(TAG), 루즈-패시브(LP) 등 4대 성향 유형으로 어떻게 분할 배치되는지 시각적 매트릭스로 규명할 수 있게 설계했습니다.

특히 이 데이터 마트는 분석용 로컬 SQLite DB와 상용 PostgreSQL DB 간에 환경 변수 값 하나로 즉각 스위칭 가능한 이중 DB 호환 커넥터를 제가 직접 설계하여 구축함으로써 서비스 이식성을 극대화했습니다."

[용어 리마인더]
- 데이터 마트: 전체 데이터웨어하우스에서 특정 분석 목적(여기서는 플레이어 성향 분류)을 위해 가공한 요약 데이터셋.
- 4대 성향 코호트:
  1. Tight-Aggressive (TAG): 강한 패로만 참여하고 들어갔을 땐 강하게 선공하는 모범적 스타일.
  2. Loose-Passive (LP): 나쁜 패로도 마구 들어가서 남의 베팅에 콜만 하며 따라다니는 초보적 스타일."
- 산점도(Scatter Plot): 2차원 평면에 점들을 찍어 두 변수(VPIP, PFR) 간의 상관관계를 시각화하는 차트.""")

    # =========================================================================
    # SLIDE 5: ZOOM-IN Slide (02. 행동 프로파일링 데이터 마트 줌인)
    # =========================================================================
    add_zoom_slide(
        prs, "02-1", "플레이어 행동 프로파일링 데이터 마트 상세", "tab2_profiling_scrolled_1780972876491.png",
        """[슬라이드 상세 대본]
"화면을 확대하여 행동 분포 차트 부분을 자세히 살펴보겠습니다.
VPIP와 PFR을 두 축으로 하는 우측 하단의 2차원 산점도를 보시면, 플레이어들의 데이터 포인트가 대각선 하단 영역에 우상향 형태로 모여 있는 강한 상관성을 확인할 수 있습니다. 

여기서 VPIP가 25%를 넘어가면서 PFR이 매우 낮은 우측 하단 영역(Loose-Passive)에 넓게 퍼져 있는 플레이어들의 데이터를 분석해 보면, 기대 손익이 아주 나쁜 적자 구렁텅이에 빠져 있음을 수학적으로 입증할 수 있습니다."
"""
    )

    # =========================================================================
    # SLIDE 6: 03. 인게임 베팅 퍼널 및 성향 코호트 분석 (Tab 3: Funnel & Cohort)
    # =========================================================================
    slide4 = create_slide_with_bg(prs)
    add_header(slide4, "03. 인게임 베팅 퍼널 및 성향 코호트 분석", "FUNNEL & COHORT ANALYSIS")
    
    add_card(
        slide4, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT, 
        "비즈니스 파산 원인 규명 및 핵심 의사결정",
        [
            "■ 베팅 진행단계별 깔때기 퍼널(Funnel) 분석",
            "  - 프리플랍(100%) ➡️ 플랍(57.0%) ➡️ 턴(44.0%) ➡️ 리버(33.0%) ➡️ 쇼다운(25.0%)",
            "  - 프리플랍에서 플랍 전환 시 43%가 탈락(Fold)하며 리스크 관리가 수행됨을 정량 증명",
            "",
            "■ 플레이 성향 코호트(Cohort)별 30세션 누적 칩 추이 추적",
            "  - Tight-Aggressive(TAG) 코호트는 장기 우상향 생존",
            "  - ★소극적 유저(Loose-Passive) 코호트는 평균 **-16,393 칩스**의 가파른 적자를 내며 파산함을 규명",
            "",
            "■ ★데이터 분석 기반의 최종 비즈니스 의사결정",
            "  - 단순히 글로 써진 공허한 플레이 가이드를 주는 방식 탈피",
            "  - **'본인 지표 변경 시 미래 자산 파산 추이를 모의실험해보는 실시간 샌드박스 시뮬레이터 기능을 서비스에 신설 및 상용 배포 결정'**"
        ]
    )
    add_image_or_placeholder(slide4, "tab3_funnel_cohort_1781002583397.png", IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    
    set_presenter_notes(slide4, """[슬라이드 4 발표 대본]
"세 번째로 이번 프로젝트에서 가장 비즈니스 임팩트가 컸던 '퍼널 및 코호트 분석(Funnel & Cohort)' 탭과 데이터 기반 의사결정 경험입니다.

먼저 포커 게임 내 베팅 라운드별로 유저가 끝까지 패를 포기하지 않고 버티는 전환율을 나타내는 '인게임 베팅 퍼널'을 설계했습니다. 프리플랍 100%에서 시작해 플랍으로 넘어갈 때 43%의 유저가 폴드(기권)하여 나쁜 패를 조기 정리하는 리스크 관리 흐름이 데이터로 입증되었습니다.

가장 중요한 분석 결과는 우측 하단 그래프인 '성향 코호트별 30세션 누적 자산 추이'입니다. 성향이 같은 플레이어들을 집단(코호트)으로 묶고, 30세션 동안 게임을 치르는 시간 경과에 따른 평균 누적 칩 자산의 변화를 추적했습니다. 분석 결과, 강한 패 위주로 공격하는 TAG 군집은 장기적으로 우상향하며 살아남았으나, 소극적으로 콜만 누르며 참견하는 Loose-Passive 군집은 평균 -16,393 칩스의 막대한 적자를 입고 급격히 파산하여 아예 서비스를 이탈하는 Churn 현상이 통계적으로 규명되었습니다.

이 차트를 기반으로 저는 최종적인 비즈니스 의사결정을 내렸습니다. '소심하게 베팅하지 말라'는 식의 텍스트 안내장을 팝업창으로 띄워봤자 유저 행동을 바꿀 수 없습니다. 대신 유저가 직접 수치 슬라이더를 조절하며 자신의 미래 자산이 우하향하며 파산하는 경로를 모의실험해볼 수 있는 '실시간 샌드박스 시뮬레이터 기능을 플랫폼에 개발하여 상용 배포하기로 결정'한 것입니다."

[용어 리마인더]
- 퍼널 분석: 유저가 서비스의 최종 도달지(여기서는 패를 오픈하는 쇼다운)까지 단계별로 진입하며 좁아지는 깔때기 모양의 이탈/생존율 분석.
- 코호트 분석: 공통 특성(여기서는 플레이 성향)을 공유하는 사용자 그룹을 나누어 시간 경과에 따른 지표 변화를 종단적으로 추적하는 기법.
- Churn(이탈): 본 분석에서는 게임 자산 파산(Bust-out)으로 인해 게임 접속을 영구 중단하는 현상.""")

    # =========================================================================
    # SLIDE 7: ZOOM-IN Slide (03. 퍼널 및 코호트 분석 줌인)
    # =========================================================================
    add_zoom_slide(
        prs, "03-1", "베팅 퍼널 및 코호트 자산 추이 상세", "tab3_funnel_cohort_scrolled_1780972892608.png",
        """[슬라이드 상세 대본]
"우측의 상세 분석 화면을 확대해 드렸습니다.
좌측의 퍼널 그래프를 보시면 각 라운드별(Pre-flop -> Flop -> Turn -> River -> Showdown)로 이탈(Fold)하는 유저들의 비율이 막대로 잘 설계되어 있습니다.

그리고 우측의 '성향 코호트별 30세션 누적 수익 추이' 시계열 차트를 보시면, 하늘색 선인 Tight-Aggressive 성향 집단만 유일하게 누적 수익이 우상향하고 있습니다. 반면에 붉은색 선인 Loose-Passive 성향 집단은 시간이 지날수록 그래프가 가파르게 무너지며 0원 선 밑으로 곤두박질치는 양상이 뚜렷하게 시각화됩니다. 이 통계적 데이터가 시뮬레이터를 개발하게 된 결정적인 비즈니스 근거입니다."
"""
    )

    # =========================================================================
    # SLIDE 8: 04. 통계적 A/B 테스트 및 머신러닝 예측 모델링 (Tab 4: A/B Test)
    # =========================================================================
    slide5 = create_slide_with_bg(prs)
    add_header(slide5, "04. 통계적 A/B 테스트 및 머신러닝 예측 모델링", "HYPOTHESIS TESTING & ML MODEL")
    
    add_card(
        slide5, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT, 
        "통계 기반 가설 검정 및 성향 예측 모델",
        [
            "■ A/B 테스트 가설 설정 및 독립표본 T-검정 수행",
            "  - 가설: 프리미엄 카드(AA/KK/QQ) 진입 시 선제공격(A)이 소극적 콜(B)보다 수익이 높을 것이다.",
            "  - T-검정 결과: T-통계량 = 0.4734, p-value = 0.6372로 나와 귀무가설 채택",
            "",
            "■ ★데이터 분석 기반의 프로세스 의사결정",
            "  - 두 전략 간의 수익 차이가 없음이 아닌, 대조군(수동 진입)의 표본수가 극히 부족(A그룹 386건 vs B그룹 64건)하여 발생한 통계적 검정력 부족 진단",
            "  - **섣부른 플레이 가이드라인 도출을 유보하고, 성향점수 매칭(PSM) 및 부트스트랩을 데이터 파이프라인에 적용해 분석 엔진을 개선하기로 의사결정**",
            "",
            "■ 머신러닝 모델 적용 및 연동",
            "  - **K-Means Clustering**: 932명 플레이어 군집 자동화",
            "  - **Random Forest Classifier**: 유저 흑자 여부 예측 (Accuracy 84%)"
        ]
    )
    add_image_or_placeholder(slide5, "tab4_ab_test_1781002593463.png", IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    
    set_presenter_notes(slide5, """[슬라이드 5 발표 대본]
"네 번째로 'A/B 테스트 가설 검정 및 머신러닝(A/B Test)' 탭과 두 번째 데이터 기반 의사결정 사례입니다.

저는 좋은 카드인 프리미엄 카드를 들었을 때 선제 레이즈로 주도권을 가져가는 A그룹의 수익이, 단순 콜만 하는 B그룹보다 유의미하게 클 것이라는 가설을 세우고 독립표본 T-검정을 수행했습니다. 결과는 우측 대시보드 검정창에 표시된 대로 p-value가 0.6372로 나와 통계적 유의성이 기각되었습니다.

여기서 저의 데이터 분석가로서의 두 번째 '프로세스적 의사결정'이 일어났습니다. 수치만 보고 섣부르게 '두 전략은 수익 차이가 없습니다'라고 발표했다면, 잘못된 결론이었을 것입니다. 저는 데이터 이면의 원인을 분석해, 수동 진입을 한 B그룹의 표본이 단 64건으로 A그룹(386건) 대비 극도로 적어 통계적 검정의 신뢰도(검정력)가 무너졌음을 찾아냈습니다.
이에 따라 저는 '섣부른 의사결정 가이드를 배포하지 않고, 성향점수 매칭(PSM) 기법과 부트스트랩 리샘플링 기술을 도입하여 표본 편향을 수학적으로 완전히 극복하는 로드맵을 수립하고 통계 엔진 개발을 보완하기로 의사결정'을 내렸습니다.

또한, K-Means 알고리즘을 통해 932명 플레이어들을 자동으로 군집 세그멘테이션하고, 유저의 초기 베팅 지표를 넣었을 때 최종 흑자/적자 여부를 84% 정확도로 맞추는 랜덤 포레스트 예측 모형을 대시보드 뒤에 탑재했습니다."

[용어 리마인더]
- 독립표본 T-검정: 두 독립적인 집단의 평균 차이가 통계적으로 의미가 있는 차이인지 검정하는 방법.
- p-value(유의확률): 귀무가설(차이가 없다)이 맞다는 전제하에 현재의 차이가 우연히 일어났을 확률. 보통 0.05 미만이어야 우연이 아닌 진짜 차이로 간주. 0.6372는 63% 확률로 우연히 발생했다는 뜻.
- K-Means 군집 분석(비지도학습): 데이터들 사이의 거리를 계산하여 정답 레이블 없이도 비슷한 특성의 데이터끼리 K개의 그룹으로 묶어주는 알고리즘.
- Random Forest(지도학습): 여러 개의 의사결정 나무 모델들을 앙상블로 투표하게 하여 정확도와 안정성을 높인 예측/분류 알고리즘.""")

    # =========================================================================
    # SLIDE 9: ZOOM-IN Slide (04. A/B 테스트 및 가설 검증 줌인)
    # =========================================================================
    add_zoom_slide(
        prs, "04-1", "A/B 테스트 및 가설 검증 결과 상세", "tab4_ab_test_bottom_1780972602234.png",
        """[슬라이드 상세 대본]
"A/B 테스트 통계 검정 결과 상세 화면을 확대해 드렸습니다.
상단의 실시간 검정 요약 카드를 보시면, 두 그룹 간의 평균 수익 분포와 신뢰구간 시각화가 나타나 있습니다.

여기서 🧪 통계 검정 지표 결과를 보시면 p-value가 0.6372가 계산되며 빨간 경고창으로 '귀무가설 채택 (A/B 간 통계적 차이가 존재하지 않음)' 문구가 나타납니다. 하지만 오차 막대의 중첩 상태와 하단 표본 수 차이를 보면, B그룹의 데이터 밀도가 극도로 떨어져 있음을 알 수 있습니다. 이 분석적 결함 진단 덕분에 데이터 수집 파이프라인의 보완이 필요함을 입증할 수 있었습니다."
"""
    )

    # =========================================================================
    # SLIDE 10: 05. 실시간 ETL 및 플레이 스타일 샌드박스 시뮬레이터 (Tab 5: Simulator)
    # =========================================================================
    slide6 = create_slide_with_bg(prs)
    add_header(slide6, "05. 실시간 ETL 및 플레이 스타일 샌드박스 시뮬레이터", "PLATFORM MaaS & SIMULATOR")
    
    add_card(
        slide6, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT, 
        "시뮬레이터 구현 및 비즈니스 행동 교정",
        [
            "■ ⚡ 실시간 데이터 적재 시뮬레이션 (Streaming ETL)",
            "  - 140개 텍스트 로그 파일의 7,993개 핸드가 DB에 적재되는 과정 모니터링",
            "  - 실시간 진행바(Progress Bar) 및 수집된 유저/핸드 수 실시간 렌더링",
            "",
            "■ 🎭 플레이 스타일 샌드박스 (Playstyle Sandbox) 기능",
            "  - 유저가 VPIP, PFR, AF 슬라이더 조절 시 실시간 성향 코호트 분류",
            "  - 머신러닝의 유클리드 거리 기준으로 DB 내 가장 지표가 유사한 실제 유저 3명을 탐색하여 자동 매칭",
            "",
            "■ ★의사결정 보조 및 비즈니스 임팩트",
            "  - 매칭 유저의 30세션 누적 자산 추이 실시간 시각 대조 렌더링",
            "  - 유저 본인의 잘못된 베팅 습관(Loose-Passive) 설정 시 가파르게 우하향하여 파산하는 모습을 직접 모의실험하게 유도해 능동적인 플레이 교정 자극"
        ]
    )
    add_image_or_placeholder(slide6, "tab5_simulator_done_1781002631767.png", IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    
    set_presenter_notes(slide6, """[슬라이드 6 발표 대본]
"다섯 번째로 본 분석 프로젝트의 최종적 성과물이자 의사결정 도구인 '시뮬레이터(Simulator)' 탭에 대해 말씀드리겠습니다.

이 화면은 크게 두 가지 기능을 탑재하고 있습니다. 상단의 '실시간 데이터 적재 시뮬레이션'은 비정형 텍스트 로그가 정제 가공을 거쳐 DB에 실시간 스트리밍 적재되는 파이프라인의 진행 과정을 보여줍니다. 버튼을 누르면 진행바가 차오르며 수집 데이터량이 실시간으로 갱신됩니다.

하단의 '플레이 스타일 샌드박스'가 바로 의사결정을 유도하는 핵심 시뮬레이터입니다. 유저가 마우스로 자신의 VPIP, PFR, AF 수치를 요리조리 움직이면, 머신러닝 거리 연산 공식을 작동시켜 DB 내의 932명 실제 플레이어들 중에서 나와 가장 전략이 똑같은 실제 유저 3명을 찾아내서 매칭합니다.
그리고 그 매칭된 유저 3명이 30세션 동안 벌어들인 누적 자산 추이를 선 그래프로 대조하여 띄워 줍니다.

만약 유저가 적자가 큰 Loose-Passive 형태로 수치를 세팅하면, 화면 우측에 매칭된 유저들의 자산이 가파르게 떨어져 결국 파산하는 결과를 즉시 확인하게 됩니다. 유저는 이 시각적 대조를 통해 자신의 잘못된 습관이 장기적으로 서비스 탈퇴(파산)를 불러올 것임을 실감하고, 최적의 범주인 VPIP 15~22%, PFR 12~20%, AF 2.0~3.5로 수치를 조정해 자산이 우상향하는 실제 유저의 포인트를 찾아 본인의 전략적 의사결정을 수정하도록 행동 유도 임팩트를 가하게 됩니다."

[용어 리마인더]
- 실시간 ETL(Extract, Transform, Load): 데이터 소스에서 원천 데이터를 가져와서 분석용으로 변형하여 데이터베이스에 적재하는 파이프라인.
- 유클리드 거리 공식: 다차원 공간 상의 두 변수 사이의 직선 거리를 계산하는 수학 공식. 여기서는 VPIP, PFR, AF 수치를 3차원 공간의 좌표로 변환하여 두 유저의 거리가 짧을수록(지표 차이가 작을수록) 닮은 꼴 유저로 간주.""")

    # =========================================================================
    # SLIDE 11: ZOOM-IN Slide (05. 시뮬레이터 줌인)
    # =========================================================================
    add_zoom_slide(
        prs, "05-1", "실시간 ETL 및 플레이 스타일 시뮬레이션 상세", "tab5_simulator_done_1781002631767.png",
        """[슬라이드 상세 대본]
"시뮬레이터 탭의 실행 화면을 전체 화면으로 확대해 드렸습니다.
좌측 하단에는 플레이어가 마우스 드래그로 조절할 수 있는 PFR, AF 등의 '플레이어 전략 성향 조절기' 슬라이더 패널이 구현되어 있습니다.

그 슬라이더를 우측으로 조정하는 즉시 우측 차트에 3개의 30세션 누적 칩 자산 선 그래프가 동적으로 리렌더링되며, 실시간으로 유사 유저들의 세션 잔존 추이를 대조 분석해 줍니다. 이처럼 정적인 데이터 보고서를 뛰어넘어 유저의 실시간 의사결정을 보조하고 게임 행동 변화를 견인하는 대화형 플랫폼을 구축하는 성과를 얻었습니다."
"""
    )

    # =========================================================================
    # SLIDE 12: 06. 프로젝트 성과 및 향후 고도화 로드맵 (Roadmap)
    # =========================================================================
    slide12 = create_slide_with_bg(prs)
    add_header(slide12, "06. 프로젝트 성과 및 향후 고도화 로드맵", "SUMMARY & TECHNICAL ROADMAP")
    
    # 2컬럼 레이아웃으로 성과와 로드맵 카드 배치 (이미지 없이 텍스트 카드 2개로 전면 구성)
    add_card(
        slide12, 0.8, 1.5, 5.6, 5.2, 
        "분석 및 엔지니어링 성과 (Successes)",
        [
            "■ 비정형 데이터의 RDBMS 정형화",
            "  - 텍스트 형태의 비정형 로그를 표준 RDBMS 스타 스키마 구조의",
            "    Fact/Dimension 테이블로 구조화하여 다차원 쿼리 조회 효율 극대화",
            "",
            "■ 분석의 서비스화 (MaaS - Model as a Service) 구현",
            "  - 단방향 결과 보고서에 그치지 않고 인터랙티브 시뮬레이션",
            "    대시보드 웹 플랫폼 서비스로 배포하여 데이터 비즈니스 가치 증명",
            "",
            "■ 듀얼 DB 이중화 아키텍처 수립",
            "  - SQLite의 IGNORE 문법과 PostgreSQL의 ON CONFLICT 문법 차이를",
            "    극복하고, 바인딩 파라미터를 추상화하여 인프라 호환성 완벽 구현"
        ]
    )
    
    add_card(
        slide12, 6.9, 1.5, 5.6, 5.2, 
        "한계점 및 향후 고도화 로드맵 (Roadmap)",
        [
            "■ 통계적 신뢰성 및 분석 고도화",
            "  - T-Test의 극심한 대조군 표본 부족 현상 및 변동성 편향 해결",
            "  - **성향점수 매칭(PSM) 및 부트스트랩 리샘플링 통계 기술**을",
            "    적재 파이프라인에 통합해 가설 검정 유의성 판단 메커니즘 보완",
            "",
            "■ 빅데이터 아키텍처 확장 계획",
            "  - 단일 RDBMS의 동시성 락 병목 및 대용량 트래픽 임계점 대응",
            "  - 수집 큐인 **Apache Kafka** 및 분산 연산 엔진인 **Apache Spark** 연동",
            "  - 구글 클라우드의 **GCP BigQuery 데이터웨어하우스**로",
            "    데이터 저장소 마이그레이션 아키텍처 설계 및 단계별 마이그레이션 착수"
        ]
    )
    
    set_presenter_notes(slide12, """[슬라이드 12 발표 대본]
"마지막으로 이번 프로젝트를 통해 얻은 주요 성과와 앞으로의 개선 로드맵에 대해 설명해 드리겠습니다.

대표적인 성과는 첫째, 텍스트 형태의 비정형 원천 로그를 데이터웨어하우스 표준 설계인 스타 스키마 형태로 RDBMS에 정형 적재하여 다차원 조회 성능을 확보했다는 점입니다. 둘째, 정적인 분석 보고서에 그치지 않고, 유저의 실제 습관 개선을 돕는 웹 시뮬레이션 서비스 플랫폼으로 완전히 구동 배포하여 데이터의 활용 가치를 극대화했다는 점입니다. 셋째, 로컬 분석용 DB와 상용 DB 간의 문법 격차를 차단하는 듀얼 아키텍처를 구현해 유연성을 입증했습니다.

동시에, 향후 보완할 한계점과 기술적 고도화 로드맵도 구체적으로 수립했습니다.
우선 A/B 테스트에서 B그룹 표본 부족으로 기각된 문제는, 표본 편향을 수학적으로 상쇄하는 성향점수 매칭(PSM)과 부트스트랩 기법을 파이프라인에 주입해 통계적 신뢰성을 확보하겠습니다.
또한 향후 수천만 건 이상의 트래픽이 집중되는 실제 프로덕션 환경에 대응하기 위해, 단일 DB의 병목을 해결할 실시간 수집 큐인 Apache Kafka, 분산 데이터 처리 엔진인 Apache Spark를 도입하고 구글 클라우드의 BigQuery DW로 마이그레이션하는 빅데이터 인프라 아키텍처를 수립하여 추진 중에 있습니다.

비정형 수집 적재부터 스타 스키마 아키텍처 설계, 비즈니스 의사결정 임팩트 증명 및 MaaS 플랫폼 배포까지 데이터 파이프라인 전체를 주도해 본 역량을 가지고, SK Planet에 입사하여 가치 있는 비즈니스 성과를 창출해 내겠습니다. 이상으로 발표를 마치겠습니다. 감사합니다."

[용어 리마인더]
- 성향점수 매칭(PSM): 대조군 표본수가 극히 적고 성향이 불균형할 때, 유사한 성향을 가진 가상의 대조군 매칭을 통해 통계적 편향을 예방하는 기법.
- 부트스트랩: 표본 데이터에서 무작위로 여러 번 중복 추출하여 신뢰구간과 표본 오차를 재산출하는 리샘플링 통계법.
- Apache Kafka: 실시간으로 대규모 데이터가 밀려 들어올 때 병목 없이 순서대로 흘려보내는 버퍼 역할을 하는 실시간 데이터 스트리밍 메세지 큐.
- Apache Spark: 여러 대의 컴퓨터로 데이터를 쪼개어 아주 빠르게 대용량 분산 연산을 수행해 주는 대표적인 빅데이터 연산 프레임워크.
- GCP BigQuery: 구글 클라우드에서 제공하는 초고속 대용량 클라우드 데이터웨어하우스 저장소 서비스.""")

    # PPT 저장
    ppt_path = "poker_analytics_presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved successfully to: {os.path.abspath(ppt_path)}")

if __name__ == "__main__":
    main()
